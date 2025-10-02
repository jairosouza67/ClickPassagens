"""
Script para gerar apresentação PowerPoint do ClickPassagens
Baseado nas imagens de referência e documentação técnica
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os
from pathlib import Path

# Configurações
DESIGN_REF_PATH = Path("design-reference")
OUTPUT_FILE = "design-mockups/ClickPassagens_Redesign_Completo.pptx"

# Cores do tema ClickPassagens
COLORS = {
    'primary_blue': RGBColor(30, 58, 138),      # #1e3a8a
    'light_blue': RGBColor(59, 130, 246),       # #3b82f6
    'gold': RGBColor(245, 158, 11),             # #f59e0b
    'light_gold': RGBColor(251, 191, 36),       # #fbbf24
    'purple': RGBColor(99, 102, 241),           # #6366f1
    'green': RGBColor(16, 185, 129),            # #10b981
    'dark_gray': RGBColor(30, 41, 59),          # #1e293b
    'medium_gray': RGBColor(100, 116, 139),     # #64748b
    'light_gray': RGBColor(241, 245, 249),      # #f1f5f9
    'white': RGBColor(255, 255, 255),           # #ffffff
}

def create_presentation():
    """Cria a apresentação completa"""
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    
    # Slide 1: Capa
    add_cover_slide(prs)
    
    # Slide 2: Índice
    add_index_slide(prs)
    
    # Slide 3: Visão Geral
    add_overview_slide(prs)
    
    # Slide 4: Paleta de Cores
    add_color_palette_slide(prs)
    
    # Slide 5-16: Telas do Sistema (com imagens de referência)
    add_screen_slides(prs)
    
    # Slide 17: Comparação Antes/Depois
    add_comparison_slide(prs)
    
    # Slide 18: Roadmap
    add_roadmap_slide(prs)
    
    # Slide 19: Tecnologias
    add_tech_stack_slide(prs)
    
    # Slide 20: Próximos Passos
    add_next_steps_slide(prs)
    
    # Salvar apresentação
    os.makedirs("design-mockups", exist_ok=True)
    prs.save(OUTPUT_FILE)
    print(f"✅ Apresentação criada com sucesso: {OUTPUT_FILE}")

def add_cover_slide(prs):
    """Slide 1: Capa"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background gradiente
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS['primary_blue']
    
    # Título principal
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = "✈️ ClickPassagens"
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(80)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Proposta de Redesign Completo\nVersão 2.0 - Outubro 2025"
    
    for paragraph in subtitle_frame.paragraphs:
        paragraph.font.size = Pt(32)
        paragraph.font.color.rgb = COLORS['light_gold']
        paragraph.alignment = PP_ALIGN.CENTER
    
    # Badge
    badge_box = slide.shapes.add_textbox(Inches(6), Inches(6.5), Inches(4), Inches(0.8))
    badge_frame = badge_box.text_frame
    badge_frame.text = "🎨 Design Premium | 📱 Mobile First | ⚡ Alta Performance"
    
    p = badge_frame.paragraphs[0]
    p.font.size = Pt(18)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

def add_index_slide(prs):
    """Slide 2: Índice"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    add_slide_title(slide, "📑 Índice da Apresentação")
    
    # Conteúdo
    content_items = [
        "1️⃣ Visão Geral do Projeto",
        "2️⃣ Paleta de Cores e Identidade Visual",
        "3️⃣ Tela 01 - Landing Page / Home",
        "4️⃣ Tela 02 - Sistema de Busca Avançada",
        "5️⃣ Tela 03 - Resultados de Busca",
        "6️⃣ Tela 04 - Filtros e Ordenação",
        "7️⃣ Tela 05 - Detalhes do Voo",
        "8️⃣ Tela 06 - Comparação de Preços (Milhas vs Dinheiro)",
        "9️⃣ Tela 07 - Orçamento Personalizado",
        "🔟 Tela 08 - Sistema de Comissões",
        "1️⃣1️⃣ Tela 09 - Planos e Assinaturas",
        "1️⃣2️⃣ Tela 10 - Cadastro e Login",
        "1️⃣3️⃣ Tela 11 - Dashboard do Usuário",
        "1️⃣4️⃣ Tela 12 - Versão Mobile (PWA)",
        "1️⃣5️⃣ Comparação: Antes vs Depois",
        "1️⃣6️⃣ Roadmap de Implementação",
        "1️⃣7️⃣ Stack Tecnológico",
        "1️⃣8️⃣ Próximos Passos"
    ]
    
    content_box = slide.shapes.add_textbox(Inches(1.5), Inches(2), Inches(13), Inches(6))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for item in content_items:
        p = text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(12)
        p.level = 0

def add_overview_slide(prs):
    """Slide 3: Visão Geral"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    add_slide_title(slide, "🎯 Visão Geral do Projeto")
    
    # Princípios de Design
    principles = [
        {"icon": "🎨", "title": "Design Moderno", "desc": "Interface premium com gradientes e animações"},
        {"icon": "⚡", "title": "Performance", "desc": "Carregamento rápido e otimização mobile"},
        {"icon": "🎯", "title": "Conversão", "desc": "CTAs estratégicos e jornada otimizada"},
        {"icon": "🔒", "title": "Confiança", "desc": "Segurança, depoimentos e badges"}
    ]
    
    x_start = 1.5
    y_start = 2.5
    box_width = 3
    box_height = 2
    spacing = 0.5
    
    for i, principle in enumerate(principles):
        col = i % 2
        row = i // 2
        x = x_start + col * (box_width + spacing)
        y = y_start + row * (box_height + spacing)
        
        # Box
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['light_gray']
        shape.line.color.rgb = COLORS['light_blue']
        shape.line.width = Pt(2)
        
        # Texto
        text_frame = shape.text_frame
        text_frame.clear()
        
        # Ícone
        p1 = text_frame.paragraphs[0]
        p1.text = principle['icon']
        p1.font.size = Pt(48)
        p1.alignment = PP_ALIGN.CENTER
        
        # Título
        p2 = text_frame.add_paragraph()
        p2.text = principle['title']
        p2.font.size = Pt(24)
        p2.font.bold = True
        p2.font.color.rgb = COLORS['primary_blue']
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(10)
        
        # Descrição
        p3 = text_frame.add_paragraph()
        p3.text = principle['desc']
        p3.font.size = Pt(14)
        p3.font.color.rgb = COLORS['medium_gray']
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(5)

def add_color_palette_slide(prs):
    """Slide 4: Paleta de Cores"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "🎨 Paleta de Cores Premium")
    
    colors_info = [
        {"name": "Azul Aviation", "rgb": COLORS['primary_blue'], "hex": "#1e3a8a", "use": "Principal"},
        {"name": "Azul Claro", "rgb": COLORS['light_blue'], "hex": "#3b82f6", "use": "Secundário"},
        {"name": "Dourado Premium", "rgb": COLORS['gold'], "hex": "#f59e0b", "use": "Destaque"},
        {"name": "Roxo Moderno", "rgb": COLORS['purple'], "hex": "#6366f1", "use": "Acento"},
        {"name": "Verde Sucesso", "rgb": COLORS['green'], "hex": "#10b981", "use": "Confirmação"},
    ]
    
    x_start = 2
    y_start = 2.5
    box_width = 2.2
    box_height = 3
    spacing = 0.3
    
    for i, color in enumerate(colors_info):
        x = x_start + i * (box_width + spacing)
        
        # Color swatch
        color_box = slide.shapes.add_shape(
            1, Inches(x), Inches(y_start), Inches(box_width), Inches(1.5)
        )
        color_box.fill.solid()
        color_box.fill.fore_color.rgb = color['rgb']
        color_box.line.width = Pt(0)
        
        # Info box
        info_box = slide.shapes.add_textbox(
            Inches(x), Inches(y_start + 1.6), Inches(box_width), Inches(1.4)
        )
        text_frame = info_box.text_frame
        text_frame.word_wrap = True
        
        # Nome
        p1 = text_frame.paragraphs[0]
        p1.text = color['name']
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLORS['dark_gray']
        p1.alignment = PP_ALIGN.CENTER
        
        # Hex
        p2 = text_frame.add_paragraph()
        p2.text = color['hex']
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['medium_gray']
        p2.alignment = PP_ALIGN.CENTER
        
        # Uso
        p3 = text_frame.add_paragraph()
        p3.text = color['use']
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLORS['medium_gray']
        p3.alignment = PP_ALIGN.CENTER

def add_screen_slides(prs):
    """Adiciona slides das telas do sistema"""
    
    screens = [
        {
            "number": "01",
            "title": "Landing Page / Home",
            "image": "Home.png",
            "features": [
                "Hero section com gradiente premium",
                "Formulário de busca em destaque",
                "Estatísticas em tempo real",
                "Logos das companhias aéreas",
                "Badges de confiança",
                "Seção 'Como Funciona'",
                "Depoimentos de clientes",
                "CTAs otimizados"
            ],
            "status": "🔄 Redesenhado"
        },
        {
            "number": "02",
            "title": "Sistema de Busca Avançada",
            "image": "Busca.png",
            "features": [
                "Autocomplete de aeroportos",
                "Seletor de datas com calendário",
                "Filtros por companhia (logos)",
                "Seleção de classe de voo",
                "Número de passageiros",
                "Toggle ida e volta / só ida",
                "Busca flexível de datas",
                "Histórico de buscas"
            ],
            "status": "✨ Melhorado"
        },
        {
            "number": "03",
            "title": "Resultados de Busca",
            "image": "Busca_efetuada.png",
            "features": [
                "Cards premium com shadows",
                "Informações hierarquizadas",
                "Preço em destaque",
                "Badge 'Melhor Oferta'",
                "Botão expandir detalhes",
                "Filtros laterais",
                "Ordenação múltipla",
                "Loading skeleton"
            ],
            "status": "🔄 Redesenhado"
        },
        {
            "number": "04",
            "title": "Filtros e Ordenação",
            "image": "Filtros.png",
            "features": [
                "Filtros por companhia",
                "Filtros por horário",
                "Filtros por número de paradas",
                "Filtros por classe",
                "Ordenação por preço",
                "Ordenação por duração",
                "Contadores de resultados",
                "Reset de filtros"
            ],
            "status": "✨ Melhorado"
        },
        {
            "number": "05",
            "title": "Visualização de Dias",
            "image": "Pesquisa_de_dias.png",
            "features": [
                "Calendário de preços",
                "Comparação de datas",
                "Preços por dia",
                "Indicação de melhor dia",
                "Navegação por mês",
                "Datas flexíveis (+/- 3 dias)",
                "Gráfico de tendências",
                "Alertas de preço"
            ],
            "status": "🆕 Nova"
        },
        {
            "number": "06",
            "title": "Comparação: Milhas vs Dinheiro",
            "image": "Busca_efetuada_2.png",
            "features": [
                "Visualização lado a lado",
                "Cálculo de economia",
                "Precificação do milheiro",
                "Destaque visual das diferenças",
                "Toggle entre visualizações",
                "Gráficos comparativos",
                "Recomendação inteligente",
                "Simulador de conversão"
            ],
            "status": "🔄 Redesenhado"
        },
        {
            "number": "07",
            "title": "Orçamento Personalizado",
            "image": "Orcamento_Personalizado.png",
            "features": [
                "Editor visual de orçamentos",
                "Upload de logo da agência",
                "Customização de cores",
                "Templates pré-definidos",
                "Geração de PDF profissional",
                "Envio automático por email",
                "Histórico de orçamentos",
                "Salvamento de configurações"
            ],
            "status": "✨ Melhorado"
        },
        {
            "number": "08",
            "title": "Sistema de Comissões",
            "image": "Comissao.png",
            "features": [
                "Configuração por companhia",
                "Comissão percentual ou fixa",
                "Comissão por passageiro/trecho",
                "Faixas de milhas",
                "Preview de cálculo",
                "Relatórios de comissões",
                "Histórico de ganhos",
                "Dashboard financeiro"
            ],
            "status": "🆕 Nova"
        },
        {
            "number": "09",
            "title": "Planos e Assinaturas",
            "image": "Planos.png",
            "features": [
                "Cards de planos premium",
                "Destaque no mais popular",
                "Comparação de recursos",
                "Toggle mensal/anual",
                "Desconto para anual",
                "Tabela comparativa",
                "FAQ sobre planos",
                "Garantia de 30 dias"
            ],
            "status": "🔄 Redesenhado"
        },
        {
            "number": "10",
            "title": "Precificação do Milheiro",
            "image": "Precos_Milheiros.png",
            "features": [
                "Configuração por companhia",
                "Preço de compra de milhas",
                "Preço de venda de milhas",
                "Margem de lucro calculada",
                "Histórico de preços",
                "Alertas de variação",
                "Sugestões de mercado",
                "Export de configurações"
            ],
            "status": "✨ Melhorado"
        },
        {
            "number": "11",
            "title": "Cadastro e Login",
            "image": "Cadastro.png",
            "features": [
                "Formulário simplificado",
                "Login social (Google, Facebook)",
                "Validação em tempo real",
                "Recuperação de senha",
                "Verificação de email",
                "Onboarding guiado",
                "Termos e privacidade",
                "Design responsivo"
            ],
            "status": "🔄 Redesenhado"
        },
        {
            "number": "12",
            "title": "Dashboard do Usuário",
            "image": "Home_2.png",
            "features": [
                "Cards com métricas",
                "Próximas viagens",
                "Histórico de buscas",
                "Perfil e preferências",
                "Alertas de preços",
                "Programa de fidelidade",
                "Notificações",
                "Documentos e faturas"
            ],
            "status": "🆕 Nova"
        }
    ]
    
    for screen in screens:
        add_screen_detail_slide(prs, screen)

def add_screen_detail_slide(prs, screen_data):
    """Adiciona slide detalhado de uma tela"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    title = f"📱 Tela {screen_data['number']} - {screen_data['title']}"
    add_slide_title(slide, title)
    
    # Status badge
    status_box = slide.shapes.add_textbox(Inches(12), Inches(0.6), Inches(3), Inches(0.5))
    status_frame = status_box.text_frame
    status_frame.text = screen_data['status']
    status_frame.paragraphs[0].font.size = Pt(16)
    status_frame.paragraphs[0].font.bold = True
    status_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Imagem (se existir)
    image_path = DESIGN_REF_PATH / screen_data['image']
    if image_path.exists():
        try:
            slide.shapes.add_picture(
                str(image_path),
                Inches(1), Inches(2),
                width=Inches(7), height=Inches(5)
            )
        except Exception as e:
            # Se falhar, adicionar placeholder
            add_image_placeholder(slide, 1, 2, 7, 5, screen_data['title'])
    else:
        add_image_placeholder(slide, 1, 2, 7, 5, screen_data['title'])
    
    # Features list
    features_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(6.5), Inches(5.5))
    text_frame = features_box.text_frame
    text_frame.word_wrap = True
    
    # Título da lista
    p_title = text_frame.paragraphs[0]
    p_title.text = "✨ Principais Funcionalidades:"
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = COLORS['primary_blue']
    p_title.space_after = Pt(10)
    
    # Features
    for feature in screen_data['features']:
        p = text_frame.add_paragraph()
        p.text = f"✓ {feature}"
        p.font.size = Pt(14)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(6)
        p.level = 0

def add_image_placeholder(slide, x, y, width, height, title):
    """Adiciona placeholder quando imagem não está disponível"""
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(x), Inches(y), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLORS['light_gray']
    shape.line.color.rgb = COLORS['medium_gray']
    shape.line.width = Pt(2)
    
    text_frame = shape.text_frame
    text_frame.clear()
    
    p1 = text_frame.paragraphs[0]
    p1.text = "🖼️"
    p1.font.size = Pt(80)
    p1.alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    p2 = text_frame.add_paragraph()
    p2.text = title
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLORS['medium_gray']
    p2.alignment = PP_ALIGN.CENTER

def add_comparison_slide(prs):
    """Slide: Comparação Antes/Depois"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "📊 Comparação: Versão Atual vs Nova Versão")
    
    comparisons = [
        ("Design Visual", "Funcional, cores básicas", "Premium com gradientes e animações"),
        ("Layout", "Single-page simples", "Multi-telas com navegação fluida"),
        ("Busca", "Formulário básico", "Autocomplete, validação e filtros avançados"),
        ("Resultados", "Lista simples", "Cards premium com badges"),
        ("Mobile", "Responsivo básico", "PWA instalável com gestos otimizados"),
        ("Performance", "Boa", "Excelente com lazy loading"),
    ]
    
    # Cabeçalho
    headers = ["Aspecto", "Versão Atual ❌", "Nova Versão ✅"]
    x_positions = [1.5, 5.5, 10]
    
    for i, header in enumerate(headers):
        box = slide.shapes.add_textbox(Inches(x_positions[i]), Inches(2), Inches(4), Inches(0.6))
        text_frame = box.text_frame
        p = text_frame.paragraphs[0]
        p.text = header
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLORS['white'] if i == 0 else COLORS['white']
        p.alignment = PP_ALIGN.CENTER
        
        # Background
        shape = slide.shapes.add_shape(
            1, Inches(x_positions[i]), Inches(2), Inches(4), Inches(0.6)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['primary_blue']
        shape.line.width = Pt(0)
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)
    
    # Linhas
    y_start = 2.8
    for i, (aspect, old, new) in enumerate(comparisons):
        y = y_start + i * 0.7
        
        # Aspecto
        box1 = slide.shapes.add_textbox(Inches(1.5), Inches(y), Inches(4), Inches(0.6))
        text_frame1 = box1.text_frame
        p1 = text_frame1.paragraphs[0]
        p1.text = aspect
        p1.font.size = Pt(14)
        p1.font.bold = True
        p1.font.color.rgb = COLORS['dark_gray']
        
        # Versão Atual
        box2 = slide.shapes.add_textbox(Inches(5.5), Inches(y), Inches(4), Inches(0.6))
        text_frame2 = box2.text_frame
        p2 = text_frame2.paragraphs[0]
        p2.text = old
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLORS['medium_gray']
        
        # Nova Versão
        box3 = slide.shapes.add_textbox(Inches(10), Inches(y), Inches(4.5), Inches(0.6))
        text_frame3 = box3.text_frame
        p3 = text_frame3.paragraphs[0]
        p3.text = new
        p3.font.size = Pt(12)
        p3.font.bold = True
        p3.font.color.rgb = COLORS['green']

def add_roadmap_slide(prs):
    """Slide: Roadmap de Implementação"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "🗓️ Roadmap de Implementação")
    
    phases = [
        {"phase": "Fase 1", "name": "Core", "items": "Landing + Busca + Resultados", "duration": "1-2 semanas"},
        {"phase": "Fase 2", "name": "Detalhes", "items": "Detalhes + Comparação + Filtros", "duration": "1 semana"},
        {"phase": "Fase 3", "name": "Checkout", "items": "Checkout + Confirmação + Pagamento", "duration": "2 semanas"},
        {"phase": "Fase 4", "name": "Usuário", "items": "Dashboard + Autenticação + Perfil", "duration": "1-2 semanas"},
        {"phase": "Fase 5", "name": "Premium", "items": "Planos + Comissões + Orçamentos", "duration": "1 semana"},
        {"phase": "Fase 6", "name": "Admin", "items": "Painel Admin + Relatórios", "duration": "1-2 semanas"},
        {"phase": "Fase 7", "name": "Mobile", "items": "PWA + Notificações Push", "duration": "1 semana"},
        {"phase": "Fase 8", "name": "Testes", "items": "QA + Performance + Deploy", "duration": "1 semana"},
    ]
    
    y_start = 2.2
    for i, phase in enumerate(phases):
        y = y_start + i * 0.65
        
        # Fase
        box1 = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(2), Inches(0.5))
        p1 = box1.text_frame.paragraphs[0]
        p1.text = f"{phase['phase']}: {phase['name']}"
        p1.font.size = Pt(16)
        p1.font.bold = True
        p1.font.color.rgb = COLORS['primary_blue']
        
        # Items
        box2 = slide.shapes.add_textbox(Inches(3.5), Inches(y), Inches(8), Inches(0.5))
        p2 = box2.text_frame.paragraphs[0]
        p2.text = phase['items']
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLORS['dark_gray']
        
        # Duration
        box3 = slide.shapes.add_textbox(Inches(12), Inches(y), Inches(3), Inches(0.5))
        p3 = box3.text_frame.paragraphs[0]
        p3.text = phase['duration']
        p3.font.size = Pt(14)
        p3.font.bold = True
        p3.font.color.rgb = COLORS['green']
        p3.alignment = PP_ALIGN.RIGHT
    
    # Total
    total_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(0.8))
    p_total = total_box.text_frame.paragraphs[0]
    p_total.text = "⏱️ Duração Total Estimada: 8-12 semanas (2-3 meses)"
    p_total.font.size = Pt(24)
    p_total.font.bold = True
    p_total.font.color.rgb = COLORS['primary_blue']
    p_total.alignment = PP_ALIGN.CENTER

def add_tech_stack_slide(prs):
    """Slide: Stack Tecnológico"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "⚙️ Stack Tecnológico")
    
    tech_categories = [
        {"icon": "⚛️", "name": "React 18", "desc": "Framework frontend moderno"},
        {"icon": "🎨", "name": "Tailwind CSS", "desc": "Utility-first CSS"},
        {"icon": "🔥", "name": "Vite", "desc": "Build tool ultra-rápido"},
        {"icon": "🐍", "name": "Flask/FastAPI", "desc": "Backend Python"},
        {"icon": "🗄️", "name": "PostgreSQL", "desc": "Banco de dados"},
        {"icon": "📱", "name": "PWA", "desc": "Progressive Web App"},
    ]
    
    x_start = 2
    y_start = 2.5
    box_width = 3.5
    box_height = 1.8
    
    for i, tech in enumerate(tech_categories):
        col = i % 3
        row = i // 3
        x = x_start + col * (box_width + 0.5)
        y = y_start + row * (box_height + 0.5)
        
        shape = slide.shapes.add_shape(
            1, Inches(x), Inches(y), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS['light_gray']
        shape.line.color.rgb = COLORS['light_blue']
        shape.line.width = Pt(2)
        
        text_frame = shape.text_frame
        text_frame.clear()
        
        p1 = text_frame.paragraphs[0]
        p1.text = tech['icon']
        p1.font.size = Pt(40)
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = text_frame.add_paragraph()
        p2.text = tech['name']
        p2.font.size = Pt(18)
        p2.font.bold = True
        p2.font.color.rgb = COLORS['primary_blue']
        p2.alignment = PP_ALIGN.CENTER
        
        p3 = text_frame.add_paragraph()
        p3.text = tech['desc']
        p3.font.size = Pt(12)
        p3.font.color.rgb = COLORS['medium_gray']
        p3.alignment = PP_ALIGN.CENTER

def add_next_steps_slide(prs):
    """Slide: Próximos Passos"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    add_slide_title(slide, "🚀 Próximos Passos")
    
    steps = [
        "1️⃣ Revisar e aprovar esta proposta de redesign",
        "2️⃣ Definir prioridades de funcionalidades",
        "3️⃣ Iniciar desenvolvimento na branch dev-melhorias",
        "4️⃣ Testes locais e ajustes",
        "5️⃣ Aprovação final e merge para master",
        "6️⃣ Deploy em produção (clickpassagens.me)",
        "7️⃣ Monitoramento e otimizações",
    ]
    
    content_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(12), Inches(4))
    text_frame = content_box.text_frame
    
    for step in steps:
        p = text_frame.add_paragraph()
        p.text = step
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['dark_gray']
        p.space_after = Pt(20)
    
    # CTA
    cta_box = slide.shapes.add_textbox(Inches(3), Inches(7), Inches(10), Inches(1))
    cta_frame = cta_box.text_frame
    cta_p = cta_frame.paragraphs[0]
    cta_p.text = "✨ Vamos transformar o ClickPassagens juntos! ✨"
    cta_p.font.size = Pt(32)
    cta_p.font.bold = True
    cta_p.font.color.rgb = COLORS['primary_blue']
    cta_p.alignment = PP_ALIGN.CENTER

def add_slide_title(slide, title_text):
    """Adiciona título padronizado ao slide"""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    
    p = title_frame.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['primary_blue']
    p.alignment = PP_ALIGN.LEFT

if __name__ == "__main__":
    print("🎨 Gerando apresentação PowerPoint do ClickPassagens...")
    print("📁 Analisando imagens de referência...")
    create_presentation()
    print(f"\n🎉 Apresentação completa criada!")
    print(f"📍 Localização: {OUTPUT_FILE}")
    print("\n✅ Pronto para apresentar!")
